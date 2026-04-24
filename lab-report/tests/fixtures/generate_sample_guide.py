"""Generate sample PPT experiment guide for testing."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_sample_guide():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "实验一：RC电路研究"
    p.font.size = Pt(44)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

    sub_box = slide1.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(0.8))
    tf2 = sub_box.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = "物理学基础实验"
    p2.font.size = Pt(24)
    p2.alignment = PP_ALIGN.CENTER

    # Slide 2: 实验目的
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    t2 = title_box2.text_frame
    t2.paragraphs[0].text = "实验目的"
    t2.paragraphs[0].font.size = Pt(36)
    t2.paragraphs[0].font.bold = True

    bullets = ["理解RC电路的基本工作原理", "掌握示波器的使用方法", "测量RC电路的充放电特性", "验证指数衰减规律"]
    y = 1.5
    for bullet in bullets:
        box = slide2.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
        tf = box.text_frame
        tf.paragraphs[0].text = f"• {bullet}"
        tf.paragraphs[0].font.size = Pt(24)
        y += 0.7

    # Slide 3: 实验原理
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    t3 = title_box3.text_frame
    t3.paragraphs[0].text = "实验原理"
    t3.paragraphs[0].font.size = Pt(36)
    t3.paragraphs[0].font.bold = True

    content_box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(2.5))
    tf3 = content_box3.text_frame
    tf3.paragraphs[0].text = "RC电路由电阻(R)和电容(C)串联组成。当电源接通时，电容开始充电，其电压随时间按指数规律上升："
    tf3.paragraphs[0].font.size = Pt(20)

    # Diagram placeholder shape
    shape = slide3.shapes.add_shape(1, Inches(3), Inches(4), Inches(4), Inches(2))  # rectangle
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(200, 200, 200)
    tf_shape = shape.text_frame
    tf_shape.paragraphs[0].text = "[电路示意图]"
    tf_shape.paragraphs[0].font.size = Pt(16)
    tf_shape.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Slide 4: 实验步骤1
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    t4 = title_box4.text_frame
    t4.paragraphs[0].text = "实验步骤 (1/4)"
    t4.paragraphs[0].font.size = Pt(36)
    t4.paragraphs[0].font.bold = True

    step_box4 = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf4 = step_box4.text_frame
    tf4.paragraphs[0].text = "1. 连接电路\n   按图连接电阻、电容和电源，确保接线牢固。\n\n2. 调节示波器\n   打开示波器，调节到合适的时基和电压档位。\n\n3. 观察波形\n   接通电源，观察电容两端的电压变化波形。"
    tf4.paragraphs[0].font.size = Pt(20)

    # Slide 5: 实验步骤2
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    t5 = title_box5.text_frame
    t5.paragraphs[0].text = "实验步骤 (2/4)"
    t5.paragraphs[0].font.size = Pt(36)
    t5.paragraphs[0].font.bold = True

    step_box5 = slide5.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf5 = step_box5.text_frame
    tf5.paragraphs[0].text = "4. 测量充电曲线\n   使用示波器 Capture 充电过程，记录电压随时间的变化。\n\n5. 测量放电曲线\n   断开电源，让电容放电，记录放电曲线。"
    tf5.paragraphs[0].font.size = Pt(20)

    # Slide 6: 实验步骤3
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    t6 = title_box6.text_frame
    t6.paragraphs[0].text = "实验步骤 (3/4)"
    t6.paragraphs[0].font.size = Pt(36)
    t6.paragraphs[0].font.bold = True

    step_box6 = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf6 = step_box6.text_frame
    tf6.paragraphs[0].text = "6. 改变参数\n   更换不同容量的电容，重复上述测量过程。\n\n7. 数据记录\n   记录不同电容下的充电和放电时间常数。"
    tf6.paragraphs[0].font.size = Pt(20)

    # Slide 7: 实验步骤4
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box7 = slide7.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    t7 = title_box7.text_frame
    t7.paragraphs[0].text = "实验步骤 (4/4)"
    t7.paragraphs[0].font.size = Pt(36)
    t7.paragraphs[0].font.bold = True

    step_box7 = slide7.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5))
    tf7 = step_box7.text_frame
    tf7.paragraphs[0].text = "8. 数据处理\n   根据测量结果计算时间常数τ=RC，与理论值比较。\n\n9. 撰写报告\n   整理实验数据，绘制V-t曲线，分析误差来源。"
    tf7.paragraphs[0].font.size = Pt(20)

    # Slide 8: 注意事项
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    title_box8 = slide8.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    t8 = title_box8.text_frame
    t8.paragraphs[0].text = "注意事项"
    t8.paragraphs[0].font.size = Pt(36)
    t8.paragraphs[0].font.bold = True

    notes = ["操作前检查电路连接是否正确", "示波器探头要正确接地", "测量时避免外界电磁干扰", "记录数据要客观真实"]
    y8 = 1.5
    for note in notes:
        box = slide8.shapes.add_textbox(Inches(1), Inches(y8), Inches(8), Inches(0.5))
        tf = box.text_frame
        tf.paragraphs[0].text = f"⚠ {note}"
        tf.paragraphs[0].font.size = Pt(22)
        y8 += 0.7

    output_path = "lab-report/tests/fixtures/sample_guide.pptx"
    prs.save(output_path)
    print(f"Saved: {output_path}")
    return output_path

if __name__ == "__main__":
    create_sample_guide()