#!/usr/bin/env python3
"""PPTX text extraction."""

import argparse
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

def parse_pptx(filepath: Path, output_format: str = "json"):
    """Parse PPTX and extract text from all slides."""
    if not HAS_PPTX:
        return {"error": "Missing dependency: python-pptx"}

    result = {
        "filename": filepath.name,
        "slide_count": 0,
        "slides": []
    }

    try:
        prs = Presentation(filepath)
        result["slide_count"] = len(prs.slides)

        for i, slide in enumerate(prs.slides):
            slide_data = {
                "number": i + 1,
                "title": "",
                "content": []
            }

            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text_frame.text.strip()
                    if text:
                        # First text shape with content is likely the title
                        if not slide_data["title"]:
                            slide_data["title"] = text
                        else:
                            slide_data["content"].append(text)

            result["slides"].append(slide_data)

    except Exception as e:
        result["error"] = str(e)

    return result

def to_markdown(result: dict) -> str:
    """Convert result to markdown format."""
    lines = [f"# {result['filename']}\n"]

    for slide in result.get("slides", []):
        lines.append(f"## Slide {slide['number']}")
        if slide.get("title"):
            lines.append(f"**{slide['title']}**\n")
        for content in slide.get("content", []):
            lines.append(content)
        lines.append("")

    return "\n".join(lines)

def main():
    parser = argparse.ArgumentParser(description='Parse PPTX file')
    parser.add_argument('--input', '-i', required=True, help='Input PPTX file')
    parser.add_argument('--format', '-f', choices=['json', 'markdown'], default='json')
    args = parser.parse_args()

    filepath = Path(args.input)
    if not filepath.exists():
        print(f"Error: File not found: {filepath}", file=sys.stderr)
        sys.exit(1)

    result = parse_pptx(filepath, args.format)

    if args.format == "markdown":
        if "error" in result:
            print(f"# Error: {result['error']}")
            sys.exit(1)
        print(to_markdown(result))
    else:
        print(json.dumps(result, indent=2, ensure_ascii=False))

    sys.exit(0 if "error" not in result else 1)

if __name__ == '__main__':
    main()